"""Build the week's results deck -> newcomb_deck.pptx (DARK theme; upload to Google Slides).

Figures (matplotlib, dark) + native pptx tables/text on black slides. Slide map (per user):
 1 headline claim (spine) + mechanism money-shot
 2 hypothesis-elimination tree + intercept-vs-slope focus box
 3 capability/disposition matrix + 4-arm intercept table
 4 K-rate(p) overlay
 5 transplant heatmap (14B)
 6 R1 iterated-game: setup + preliminary (run in progress)
 7 mechanism-credibility: abstract vs binding prompt + m0->m3 table
 8 caveats / confounders / further work
"""
from __future__ import annotations
import os
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

FIG = "results/deck"; os.makedirs(FIG, exist_ok=True)

# ---- dark palette ----
BG_HEX, FG_HEX, MUT_HEX, ACC_HEX = "#0A0A0A", "#E8E8E8", "#9AA4AF", "#4F9DE0"
BG   = RGBColor(0x0A, 0x0A, 0x0A); FG = RGBColor(0xE8, 0xE8, 0xE8)
MUT  = RGBColor(0x9A, 0xA4, 0xAF); ACC = RGBColor(0x4F, 0x9D, 0xE0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF); CELL = RGBColor(0x16, 0x1A, 0x21); BOX = RGBColor(0x14, 0x20, 0x2E)
# card fills for diagrams (dark-tinted; light text reads on them)
C_NEUT, C_GREEN, C_RED, C_AMBER, C_BLUE = "#1b2430", "#163a20", "#3a1a1d", "#3a2c12", "#16263a"
EDGE = "#5a6675"

plt.rcParams.update({
    "figure.facecolor": BG_HEX, "axes.facecolor": BG_HEX, "savefig.facecolor": BG_HEX,
    "text.color": FG_HEX, "axes.labelcolor": FG_HEX, "axes.titlecolor": FG_HEX,
    "axes.edgecolor": "#666", "xtick.color": FG_HEX, "ytick.color": FG_HEX, "grid.color": "#333",
})

# ---------------------------------------------------------------- figures
def _box(ax, x, y, w, h, text, fc, fs=11, tc=FG_HEX):
    ax.add_patch(FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.02,rounding_size=0.08",
                                fc=fc, ec=EDGE, lw=1.4))
    ax.text(x + w / 2, y + h / 2, text, ha="center", va="center", fontsize=fs, color=tc, wrap=True)

def _save(fig, path):
    fig.savefig(path, dpi=150, bbox_inches="tight", facecolor=BG_HEX); plt.close(fig)

def fig_moneyshot(path):
    fig, ax = plt.subplots(figsize=(12.4, 2.5)); ax.axis("off"); ax.set_xlim(0, 12.4); ax.set_ylim(0, 2.5)
    xs = [0.1, 3.35, 6.75, 10.3]; w = [3.0, 3.15, 3.3, 2.0]; h = 1.7; y = 0.45
    txt = ["Newcomb prompt\n(stated accuracy p)",
           "Model REPRESENTS the EV /\nstates P(full|action)=2p−1\ncorrectly",
           "Commitment step:\nDOMINANCE pull\n(“don't forgo the\nguaranteed box”)\nOVERRIDES",
           "→ two-box\n(CDT)"]
    fc = [C_NEUT, C_GREEN, C_RED, C_AMBER]
    for x, ww, t, c in zip(xs, w, txt, fc):
        _box(ax, x, y, ww, h, t, c, fs=11.5)
    for i in range(3):
        ax.add_patch(FancyArrowPatch((xs[i] + w[i], y + h / 2), (xs[i + 1], y + h / 2),
                     arrowstyle="-|>", mutation_scale=18, lw=1.7, color="#bbb"))
    ax.text(6.2, 2.35, "the override happens AFTER correct representation — that is the whole claim",
            ha="center", fontsize=10.5, style="italic", color=MUT_HEX)
    _save(fig, path)

def fig_tree(path):
    fig, ax = plt.subplots(figsize=(9.6, 6.6)); ax.axis("off"); ax.set_xlim(0, 9.6); ax.set_ylim(0, 6.6)
    nodes = [
        ("Flat K-rate(p): the model doesn't track p", C_BLUE, ""),
        ("Just a format limit? (single-token)", C_NEUT, "✗  CoT pulls toward CDT; apparent slope was a\n      scoring artifact → flat (de-noised)"),
        ("Can't compute the EV?", C_NEUT, "✗  refuses even when handed the EV"),
        ("Fixed by scale? (14B)", C_NEUT, "✗  worse — more committed CDT"),
        ("A framing artifact? (predictor credibility)", C_NEUT, "~  partly (modest, ≤10 pts)"),
        ("Dissolved by reasoning? (R1-Distill)", C_GREEN, "✓  tracks p, crossover at p*"),
    ]
    y = 6.0; bh = 0.62; gap = 0.34; bx = 0.2; bw = 5.4; cy = []
    for t, c, v in nodes:
        _box(ax, bx, y, bw, bh, t, c, fs=11)
        if v:
            col = "#5fd37a" if v.startswith("✓") else ("#e0a64f" if v.startswith("~") else "#e06a6a")
            ax.text(bx + bw + 0.25, y + bh / 2, v, ha="left", va="center", fontsize=10.3, color=col)
        cy.append(y); y -= (bh + gap)
    for i in range(len(nodes) - 1):
        ax.add_patch(FancyArrowPatch((bx + bw / 2, cy[i]), (bx + bw / 2, cy[i + 1] + bh),
                     arrowstyle="-|>", mutation_scale=14, lw=1.4, color="#888"))
    _box(ax, 0.2, y - 0.15, 9.0, 0.95,
         "⇒ a reasoning-dissolvable DISPOSITION at the commitment step\n"
         "(not a capability gap) — next: can an environment install the rule?  R1 iterated game (running)",
         C_BLUE, fs=11)
    _save(fig, path)

def fig_krate(path):
    fig, ax = plt.subplots(figsize=(8.2, 5.0))
    bp = [0.5, 0.6, 0.7, 0.75, 0.8, 0.85, 0.9, 0.99]
    base = [0.90, 0.95, 0.90, 0.80, 0.70, 0.80, 0.70, 0.90]
    caus = [0.0] * 8; evid = [0.50, 0.45, 0.50, 0.40, 0.50, 0.50, 0.45, 0.35]
    r1p = [0.5, 0.7, 0.8, 0.9, 0.99]; r1 = [0.08, 0.50, 0.58, 0.83, 0.75]
    ax.plot(bp, base, "o--", color="#4F9DE0", lw=2, label="3B base — flat-high (reflex)")
    ax.plot(bp, evid, "s--", color="#E0A64F", lw=2, label="3B evidential-RL — flat ~0.5")
    ax.plot(bp, caus, "v--", color="#E06A6A", lw=2, label="3B causal-RL — flat-zero")
    ax.plot(r1p, r1, "D-", color="#5FD37A", lw=2.6, ms=8, label="R1-8B (reasoning) — TRACKS p")
    ax.axvline(0.8, ls=":", color="#888"); ax.text(0.805, 0.04, "p*=0.8", color="#aaa", fontsize=9)
    ax.set_xlabel("stated predictor accuracy p"); ax.set_ylabel("K-rate (one-box / non-CDT rate)")
    ax.set_ylim(-0.05, 1.05); ax.set_title("Everything is flat in p — except the reasoning model")
    leg = ax.legend(fontsize=9, loc="center left", facecolor="#14181E", edgecolor="#555", labelcolor=FG_HEX)
    ax.grid(alpha=0.3)
    _save(fig, path)

def fig_heatmap(path):
    import numpy as np
    conds = ["none", "variables", "formulas", "numeric_evs", "comparison", "full"]
    ps = [0.5, 0.6, 0.7, 0.75, 0.85, 0.9, 0.99]
    M = {"none": [.527, .382, .590, .471, .431, .427, .550],
         "variables": [.711, .652, .774, .747, .179, .146, .174],
         "formulas": [.924, .895, .927, .920, .176, .047, .002],
         "numeric_evs": [1.0, .629, 1.0, 1.0, .373, .316, .000],
         "comparison": [.923, .833, .924, .917, .667, .956, .585],
         "full": [1.0, 1.0, 1.0, 1.0, .666, .779, .146]}
    A = np.array([M[c] for c in conds])
    fig, ax = plt.subplots(figsize=(8.4, 4.8))
    im = ax.imshow(A, cmap="RdYlGn", vmin=0, vmax=1, aspect="auto")
    ax.set_xticks(range(len(ps))); ax.set_xticklabels(ps); ax.set_yticks(range(len(conds))); ax.set_yticklabels(conds)
    for i in range(len(conds)):
        for j in range(len(ps)):
            ax.text(j, i, f"{A[i,j]:.2f}", ha="center", va="center", fontsize=9,
                    color="black" if 0.25 < A[i, j] < 0.8 else "white")
    ax.axvline(3.5, color="white", lw=3)   # p*=0.8 divider (between p=0.75 and p=0.85)
    ax.set_xlabel("stated p", labelpad=8); ax.set_ylabel("EV aid supplied")
    ax.text(0.27, -0.22, "← below p* = 0.8\nEV-optimal = TWO-box\n(the guaranteed-box move)",
            transform=ax.transAxes, ha="center", va="top", color="#E8E8E8", fontsize=9, fontweight="bold")
    ax.text(0.79, -0.22, "above p* = 0.8 →\nEV-optimal = ONE-box\n(forgo the guaranteed box)",
            transform=ax.transAxes, ha="center", va="top", color="#E8E8E8", fontsize=9, fontweight="bold")
    ax.set_title("14B handed the EV: P(EV-optimal action)\nrefuses to one-box at high p even with the full calc")
    cb = fig.colorbar(im, ax=ax, fraction=0.046)
    cb.set_label("P(optimal)", color=FG_HEX); cb.ax.yaxis.set_tick_params(color=FG_HEX)
    plt.setp(cb.ax.get_yticklabels(), color=FG_HEX); cb.outline.set_edgecolor("#666")
    _save(fig, path)

def fig_r1game(path):
    fig, ax = plt.subplots(figsize=(9.4, 4.0)); ax.axis("off"); ax.set_xlim(0, 9.4); ax.set_ylim(0, 4.0)
    _box(ax, 0.3, 1.4, 3.0, 1.3, "POLICY (R1)\nreasons → acts\n(one-box / two-box)", C_GREEN, 11)
    _box(ax, 6.1, 1.4, 3.0, 1.3, "PREDICTOR =\nlagged snapshot of policy\nREASONS → predicts", C_BLUE, 11)
    ax.add_patch(FancyArrowPatch((3.3, 2.3), (6.1, 2.3), arrowstyle="-|>", mutation_scale=16, lw=1.6, color="#bbb"))
    ax.text(4.7, 2.55, "action", ha="center", fontsize=9.5, color=FG_HEX)
    ax.add_patch(FancyArrowPatch((6.1, 1.7), (3.3, 1.7), arrowstyle="-|>", mutation_scale=16, lw=1.6, color="#bbb"))
    ax.text(4.7, 1.32, "fills box w/ accuracy p_eff = P_pred(one-box)", ha="center", fontsize=9.5, color=FG_HEX)
    ax.text(4.7, 3.6, "Because the predictor REASONS, p_eff tracks stated p → reward conditional from step 0",
            ha="center", fontsize=10.5, style="italic", color=MUT_HEX)
    ax.text(4.7, 0.5, "Q: is the conditional rule (one-box high p / two-box low p) a self-consistent fixed point?",
            ha="center", fontsize=10.5, color="#4F9DE0")
    _save(fig, path)


def fig_recipe(path):
    import numpy as np
    fig, (axL, axR) = plt.subplots(1, 2, figsize=(12.2, 4.6), gridspec_kw={"width_ratios": [1.3, 1]})
    axL.axis("off"); axL.set_xlim(0, 10); axL.set_ylim(0, 10)
    axL.text(5, 9.6, "A predictor (accuracy p) has already forecast your choice and\nfilled the opaque box iff it predicted ONE-box.  Then you pick:",
             ha="center", va="top", fontsize=10.5, color=FG_HEX)
    _box(axL, 3.4, 6.9, 3.2, 1.0, "you choose", C_NEUT, 12)
    _box(axL, 0.2, 4.0, 4.3, 1.5, "ONE-box\n(opaque box only)\nEV = p\u00B7B", C_GREEN, 11)
    _box(axL, 5.5, 4.0, 4.3, 1.5, "TWO-box\n(opaque + guaranteed S)\nEV = S + (1\u2212p)\u00B7B", C_AMBER, 11)
    axL.add_patch(FancyArrowPatch((4.4, 6.9), (2.4, 5.5), arrowstyle="-|>", mutation_scale=14, lw=1.5, color="#888"))
    axL.add_patch(FancyArrowPatch((5.6, 6.9), (7.6, 5.5), arrowstyle="-|>", mutation_scale=14, lw=1.5, color="#888"))
    axL.text(5, 3.1, "one-box optimal  \u21D4  p > p* = (S+B)/(2B) = 0.8     [B=100, S=60]",
             ha="center", fontsize=11.5, color="#4F9DE0")
    axL.text(5, 1.6, "CDT: boxes already fixed \u2192 two-boxing dominates (+S for any p).\n"
             "EDT: your choice is evidence about the fill \u2192 condition on p.",
             ha="center", va="top", fontsize=10, color=MUT_HEX)
    p = np.linspace(0.5, 1.0, 100)
    axR.plot(p, [0.8] * len(p), color="#E0A64F", lw=2.5, label="recitation: FLAT in p")
    axR.plot(p, np.where(p < 0.8, 0.12, 0.9), color="#5FD37A", lw=2.5, label="reasoner: STEP at p*")
    axR.axvline(0.8, ls=":", color="#888"); axR.text(0.807, 0.02, "p*", color="#aaa", fontsize=9)
    axR.set_xlabel("stated p"); axR.set_ylabel("K-rate (one-box rate)"); axR.set_ylim(-0.05, 1.05)
    axR.set_title("the signature we test for", fontsize=11)
    axR.legend(fontsize=9, loc="center left", facecolor="#14181E", edgecolor="#555", labelcolor=FG_HEX)
    axR.grid(alpha=0.3)
    _save(fig, path)

def fig_intercept(path):
    import numpy as np
    fig, ax = plt.subplots(figsize=(4.7, 4.5))
    p = np.linspace(0.5, 1.0, 50)
    ax.plot(p, [0.8] * 50, color="#4F9DE0", lw=2.5); ax.text(0.51, 0.84, "base", color="#4F9DE0", fontsize=9)
    ax.plot(p, [0.45] * 50, color="#E0A64F", lw=2.5); ax.text(0.51, 0.49, "evidential", color="#E0A64F", fontsize=9)
    ax.plot(p, [0.02] * 50, color="#E06A6A", lw=2.5); ax.text(0.51, 0.07, "causal", color="#E06A6A", fontsize=9)
    ax.plot(p, np.where(p < 0.8, 0.2, 0.9), ls="--", color="#5FD37A", lw=2.2)
    ax.text(0.6, 0.66, "the slope\nRL never installs", color="#5FD37A", fontsize=9)
    ax.annotate("", xy=(0.535, 0.8), xytext=(0.535, 0.02), arrowprops=dict(arrowstyle="<->", color="#aaa", lw=1.4))
    ax.text(0.55, 0.4, "RL moves\nthe LEVEL", color="#ccc", fontsize=9)
    ax.axvline(0.8, ls=":", color="#888")
    ax.set_xlabel("stated p"); ax.set_ylabel("K-rate"); ax.set_ylim(-0.05, 1.05)
    ax.set_title("RL moves the intercept, not the slope", fontsize=10.5); ax.grid(alpha=0.3)
    _save(fig, path)

for fn, p in [(fig_recipe, "recipe"), (fig_intercept, "intercept"), (fig_moneyshot, "moneyshot"), (fig_tree, "tree"), (fig_krate, "krate"),
              (fig_heatmap, "heatmap"), (fig_r1game, "r1game")]:
    fn(f"{FIG}/{p}.png")
print("figures done")

# ---------------------------------------------------------------- pptx
prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def slide(title, sub=None):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    tb = s.shapes.add_textbox(Inches(0.45), Inches(0.22), Inches(12.4), Inches(0.95)).text_frame
    tb.word_wrap = True; p = tb.paragraphs[0]; p.text = title
    p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = FG
    if sub:
        sp = tb.add_paragraph(); sp.text = sub; sp.font.size = Pt(13); sp.font.italic = True; sp.font.color.rgb = MUT
    return s

def pic(s, path, x, y, w):
    s.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))

def bullets(s, items, x, y, w, h, size=15):
    tf = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame; tf.word_wrap = True
    for i, (lvl, t) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("• " if lvl == 0 else "    – ") + t
        p.font.size = Pt(size if lvl == 0 else size - 2); p.font.color.rgb = FG if lvl == 0 else MUT
        p.space_after = Pt(4)

def table(s, rows, x, y, w, h, fs=12, colw=None):
    nr, nc = len(rows), len(rows[0])
    t = s.shapes.add_table(nr, nc, Inches(x), Inches(y), Inches(w), Inches(h)).table
    t.first_row = False; t.horz_banding = False
    if colw:
        for j, cw in enumerate(colw):
            t.columns[j].width = Inches(cw)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = t.cell(i, j); c.text = str(val); c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.fill.solid(); c.fill.fore_color.rgb = ACC if i == 0 else CELL
            para = c.text_frame.paragraphs[0]; para.font.size = Pt(fs)
            para.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            para.font.bold = (i == 0); para.font.color.rgb = WHITE if i == 0 else FG
    return t

def footnote(s, text, y=6.95, x=0.5, w=12.3, size=10):
    tf = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.5)).text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(size); p.font.italic = True; p.font.color.rgb = MUT

def focusbox(s, text, x, y, w, h):
    sh = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tf = sh.text_frame; tf.word_wrap = True
    sh.fill.solid(); sh.fill.fore_color.rgb = BOX; sh.line.color.rgb = ACC; sh.line.width = Pt(1.5)
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(13.5); p.font.color.rgb = FG

# S0 — Newcomb recipe (orient the listener)
s = slide("Newcomb's problem \u2014 the recipe", "how to read every chart that follows")
pic(s, f"{FIG}/recipe.png", 0.55, 1.45, 12.2)
bullets(s, [(0, "We sweep the stated accuracy p and watch the one-box (K) rate: a flat line = recitation/"
            "disposition; a step at p* = the model is actually doing the expected-value reasoning.")], 0.6, 6.85, 12.2, 0.5, size=12)

# S1 — headline + money-shot
s = slide("The claim", "Newcomb decision-theory in small/mid LLMs — week-in-review")
fb = s.shapes.add_textbox(Inches(0.5), Inches(1.25), Inches(12.3), Inches(1.7)).text_frame; fb.word_wrap = True
p = fb.paragraphs[0]
p.text = ("A small/mid LLM doesn't fail Newcomb because it can't do the math — it fails because of a "
          "DISPOSITION at the action-commitment step that overrides explicit computation, sharpens with "
          "scale, and is dissolved only by test-time reasoning.")
p.font.size = Pt(19); p.font.bold = True; p.font.color.rgb = FG
p2 = fb.add_paragraph(); p2.text = ("Competence (slope) and disposition (intercept) are orthogonal and "
          "separately movable: RL moves the lean, not the rule."); p2.font.size = Pt(14); p2.font.color.rgb = MUT
pic(s, f"{FIG}/moneyshot.png", 0.5, 3.4, 12.3)

# S2 — hypothesis tree + focus box
s = slide("How we got there — every alternative explanation, tested and killed")
pic(s, f"{FIG}/tree.png", 0.3, 1.2, 8.4)
focusbox(s, ("INTERCEPT vs SLOPE.\n\nAcross every RL arm the model's overall LEAN (intercept) moves "
             "freely — any direction, even against its prior — but the p-conditional RULE (slope) never "
             "forms.\n\nRL elicits/reweights latent disposition; it does not install competence.\n\n"
             "→ confirmed at the logit level, per-item, and under a fair conditioning-only objective."),
         8.95, 1.4, 4.05, 5.4)

# S3 — capability matrix + 4-arm
s = slide("Disposition ≠ capability — two receipts")
bullets(s, [(0, "Capability / disposition matrix (the model can; it won't):")], 0.5, 1.15, 12, 0.4, size=14)
table(s, [
    ["model", "no-think reflex", "tracks p when reasoning?", "1-box @ p=.99 given full EV", "verdict"],
    ["3B", "~0.83 one-box", "no (flat, de-noised)", "~0.53", "baseline"],
    ["14B", "~0.50", "weak (+0.09, n.s.)", "0–15%  (worse!)", "scale ⇒ more CDT"],
    ["R1-8B", "1.00 one-box", "YES (slope, crossover p*)", "— follows EV in CoT", "reasoning ⇒ EV-rational"],
], 0.5, 1.6, 12.3, 1.6, fs=12.5, colw=[1.2, 2.2, 3.0, 3.2, 2.7])
bullets(s, [(0, "RL moves only the intercept (logit level) — four arms, four levels, slope ≈ 0:")], 0.5, 3.7, 12, 0.4, size=14)
table(s, [
    ["RL arm", "P(non-CDT)", "logit margin", "slope in p", "what RL did"],
    ["base", "0.805", "+3.46", "≈ 0", "—"],
    ["causal", "0.000", "−18.16", "≈ 0", "saturated two-box"],
    ["evidential", "0.457", "−0.51", "≈ 0", "pushed to indifference"],
    ["modelpred", "0.718", "+8.36", "≈ 0", "sharpened one-box"],
], 0.5, 4.15, 12.3, 2.0, fs=12.5, colw=[2.0, 2.2, 2.2, 1.8, 4.1])

# S4 — K-rate overlay
s = slide("RL moves the level, never the shape — only reasoning bends the curve")
pic(s, f"{FIG}/krate.png", 2.4, 1.35, 8.5)
bullets(s, [(0, "Flat in p = recitation/disposition; a step at p* = structural EV reasoning. "
             "Only R1 (reasoning) tracks p, with the crossover at the theoretical p*=0.8.")], 0.6, 6.7, 12, 0.6, size=13)

# S5 — transplant heatmap
s = slide("“We handed it the answer and it still wouldn't act on it”")
pic(s, f"{FIG}/heatmap.png", 2.5, 1.35, 8.3)
bullets(s, [(0, "Supplying the EV raises P(optimal) at LOW p (where EV agrees with grabbing the guaranteed box) "
             "but FLIPS it negative at HIGH p — the 14B one-boxes 0–15% at p=0.99 even told “EV(one-box) "
             "> EV(two-box)”. Spelling out the math is counter-productive: it surfaces the guaranteed reward.")],
        0.6, 6.7, 12.1, 0.6, size=12.5)

# S6 — R1 iterated game
s = slide("Does the iterated game INSTALL the rule? — preliminary: no (and confounded)",
          "R1 self-snapshot, predictor = a lagging REASONING copy of the policy; 3 runs × 30 steps")
pic(s, "results/r1_calib_dynamics.png", 0.3, 1.45, 8.1)
bullets(s, [
    (0, "Setup: the predictor (lagging snapshot) ALSO reasons → p_eff tracks p (0b confirmed, slope +0.50)."),
    (0, "Result: NO stable conditional fixed point. The loop drifts into the SAME flat self-fulfilling basins as the 3B:"),
    (1, "seed1 → one-box-flat (slope 0); seed0 → two-box-ward. Seed/noise picks the basin."),
    (0, "The calib +1.0 @ step 10 was a TRANSIENT — slopes wobble, never settle. kl=0 ≈ kl=.02 (not a KL artifact)."),
    (0, "⚠ CONFOUND: predictor closed </think> only 13–41% at the 2048 budget → p_eff on truncated chains → inconclusive."),
    (1, "next: re-run at ~3584–4096 tok so the predictor finishes; hysteresis (committed seeds) now motivated."),
], 8.55, 1.4, 4.65, 5.5, size=11.5)

# S7 — mechanism credibility
s = slide("Was the CDT lean just a prompt artifact? — partly, but modest")
bullets(s, [
    (0, "Abstract (m0): “the process correctly identifies the choices of agents like you X% of the time” "
        "— a population statistic a causal reasoner is RIGHT to treat as non-binding."),
    (0, "Binding (m3): “…executes an exact copy of your own decision procedure on this very problem…” "
        "— the honest framing for a self-snapshot predictor."),
], 0.5, 1.2, 12.3, 1.6, size=13.5)
table(s, [
    ["predictor framing", "one-box rate", "credence gap_adj (→p-tracking)"],
    ["m0  statistical (abstract)", "0.59", "0.21"],
    ["m0pad  length placebo", "0.59", "0.17"],
    ["m1  individual model", "0.61", "0.19"],
    ["m2  process-scan", "0.63", "0.30"],
    ["m3  exact-copy (binding)", "0.69", "0.31"],
], 0.5, 3.05, 9.5, 2.4, fs=13, colw=[4.5, 2.5, 2.5])
focusbox(s, ("Credibility lifts EDT m0→m3 (placebo-controlled) — but only ~10 pts.\n\nDominance "
             "SURVIVES the most binding predictor (softened, not dissolved). Framing explains a slice; "
             "the disposition is robust to credibility."), 10.2, 3.05, 2.8, 2.4)

# S8 — caveats / further work
s = slide("Caveats, confounders & next experiments")
bullets(s, [
    (0, "CAVEATS / CONFOUNDERS"),
    (1, "Small/mid models; several single-run n≈12–20 cells (14B especially)."),
    (1, "14B>3B-CDT direction suggestive, mechanism open (sharper deliberation? tuning? noise?)."),
    (1, "Abstract-prompt confound real but modest (≤10 pts)."),
    (1, "“reflex-EDT / CoT-CDT” is a small-model fact: reflex drifts off EDT with scale, and R1 "
        "*reasoning* ≠ 3B *CoT* (different operations) — “deliberation→CDT” is not scale-invariant."),
    (1, "Discipline: many apparent effects dissolved under de-noising (+0.50 CoT slope, +0.16 SFT slope, "
        "causal-flip = KL artifact). Reported claims are the survivors."),
    (0, "FURTHER WORK"),
    (1, "R1 self-snapshot iterated game (running) — does self-reference install/stabilize the conditional rule?"),
    (1, "Seeded in-family CDT ladder (0.5B→32B) — confirm/explain 14B>3B; cuts against capability→EDT."),
    (1, "Reflex × CoT × native-reasoning matrix; anti-Newcomb camouflage (Newcomb-story vs genuine EV-action failure)."),
], 0.5, 1.2, 12.4, 5.8, size=14)

out = "newcomb_deck.pptx"; prs.save(out)
print(f"wrote {out}  ({len(prs.slides)} slides)")
