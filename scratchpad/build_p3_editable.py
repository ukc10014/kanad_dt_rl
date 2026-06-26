"""Rebuild deck p.3 ("How we got there ...") with the hypothesis tree as NATIVE, UNGROUPED
pptx shapes (rounded rects + arrow connectors + text boxes) instead of the tree.png image — so
every element stays editable after import to Google Slides. Styling mirrors build_deck.py.

Output: newcomb_deck_p3_editable.pptx  (single slide)
"""
from __future__ import annotations
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ---- dark palette (from build_deck.py) ----
BG  = RGBColor(0x0A, 0x0A, 0x0A); FG = RGBColor(0xE8, 0xE8, 0xE8); MUT = RGBColor(0x9A, 0xA4, 0xAF)
ACC = RGBColor(0x4F, 0x9D, 0xE0); BOX = RGBColor(0x14, 0x20, 0x2E)

def hx(s): return RGBColor(int(s[1:3], 16), int(s[3:5], 16), int(s[5:7], 16))
C_NEUT, C_GREEN, C_BLUE = hx("#1b2430"), hx("#163a20"), hx("#16263a")
EDGE = hx("#5a6675")
V_GREEN, V_AMBER, V_RED = hx("#5fd37a"), hx("#e0a64f"), hx("#e06a6a")

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
s = prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid(); s.background.fill.fore_color.rgb = BG

# ---- title (native textbox, as in build_deck.slide) ----
tb = s.shapes.add_textbox(Inches(0.45), Inches(0.22), Inches(12.4), Inches(0.95)).text_frame
tb.word_wrap = True; p = tb.paragraphs[0]
p.text = "How we got there — every alternative explanation, tested and killed"
p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = FG

# ---- helpers ----
def node(x, y, w, h, text, fill, fs=11, bold=False):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = EDGE; sh.line.width = Pt(1.4)
    sh.shadow.inherit = False
    tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    pp = tf.paragraphs[0]; pp.text = text; pp.alignment = PP_ALIGN.CENTER
    pp.font.size = Pt(fs); pp.font.color.rgb = FG; pp.font.bold = bold
    return sh

def verdict(x, y, w, h, text, color, fs=10.3):
    tf = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
    tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    pp = tf.paragraphs[0]; pp.text = text; pp.font.size = Pt(fs); pp.font.color.rgb = color
    return tf

def arrow(x, y0, y1):
    """vertical down arrow connector from (x,y0) to (x,y1)."""
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y0), Inches(x), Inches(y1))
    c.line.color.rgb = hx("#888888"); c.line.width = Pt(1.4)
    ln = c.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    return c

# ---- tree layout (region roughly matches the old image: x 0.35–8.7, y 1.3–6.9) ----
nodes = [
    ("Flat K-rate(p): the model doesn't track p",          C_BLUE,  None,    None),
    ("Just a format limit?  (single-token)",               C_NEUT,  "✗  CoT pulls toward CDT; apparent slope\n     was a scoring artifact → flat (de-noised)", V_RED),
    ("Can't compute the EV?",                              C_NEUT,  "✗  refuses even when handed the EV",        V_RED),
    ("Fixed by scale?  (14B)",                            C_NEUT,  "✗  worse — more committed CDT",             V_RED),
    ("A framing artifact?  (predictor credibility)",       C_NEUT,  "~  partly (modest, ≤10 pts)",              V_AMBER),
    ("Dissolved by reasoning?  (R1-Distill)",              C_GREEN, "✓  tracks p, crossover at p*",             V_GREEN),
]

bx, bw, bh = 0.35, 5.0, 0.6
gap = 0.22
vx = bx + 0.30                      # verdict text x  (slightly right of node)
top = 1.35
cx = bx + bw / 2                    # arrow column (node horizontal center)
ys = []
for i, (t, fill, v, vc) in enumerate(nodes):
    y = top + i * (bh + gap); ys.append(y)
    node(bx, y, bw, bh, t, fill, fs=11, bold=(i == 0))
    if v:
        verdict(bx + bw + 0.20, y - 0.02, 3.05, bh + 0.04, v, vc)

# arrows between consecutive nodes
for i in range(len(nodes) - 1):
    arrow(cx, ys[i] + bh, ys[i + 1])

# conclusion box (full width), with arrow into it
concl_y = ys[-1] + bh + gap
arrow(cx, ys[-1] + bh, concl_y)
cb = node(bx, concl_y, 8.35, 0.66,
          "⇒ a reasoning-dissolvable DISPOSITION at the commitment step (not a capability gap)\n"
          "next: can an environment install the rule?  R1 iterated game (running)",
          C_BLUE, fs=11)

# ---- focus box (already native in the original deck) ----
fbx = s.shapes.add_textbox(Inches(8.95), Inches(1.4), Inches(4.05), Inches(5.4))
fbx.fill.solid(); fbx.fill.fore_color.rgb = BOX; fbx.line.color.rgb = ACC; fbx.line.width = Pt(1.5)
ftf = fbx.text_frame; ftf.word_wrap = True
fp = ftf.paragraphs[0]
fp.text = ("INTERCEPT vs SLOPE.\n\nAcross every RL arm the model's overall LEAN (intercept) moves "
           "freely — any direction, even against its prior — but the p-conditional RULE (slope) never "
           "forms.\n\nRL elicits/reweights latent disposition; it does not install competence.\n\n"
           "→ confirmed at the logit level, per-item, and under a fair conditioning-only objective.")
fp.font.size = Pt(13.5); fp.font.color.rgb = FG

out = "newcomb_deck_p3_editable.pptx"
prs.save(out)
print(f"wrote {out}  ({len(prs.slides)} slide, native shapes — no group, no image)")
